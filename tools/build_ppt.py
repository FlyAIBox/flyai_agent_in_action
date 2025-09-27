#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Generate PPTX decks from docs/syllabus.yaml.

Usage examples:
  python tools/build_ppt.py                      # build all lessons into slides/
  python tools/build_ppt.py --chapter 1          # only chapter 1
  python tools/build_ppt.py --grouping chapter   # one deck per chapter
  python tools/build_ppt.py --filter 13 14 15    # only specific lesson numbers

Requires:
  pip install -r tools/requirements-slides.txt
"""

import argparse
import os
import re
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def load_syllabus(path: Path) -> Dict[str, Any]:
    with path.open('r', encoding='utf-8') as f:
        data = yaml.safe_load(f)
    if not isinstance(data, dict) or 'syllabus' not in data:
        raise ValueError('Invalid syllabus yaml: missing root "syllabus"')
    return data


def safe_filename(name: str) -> str:
    s = re.sub(r"[^0-9A-Za-z\u4e00-\u9fa5._-]+", "_", name)
    s = s.strip("._-")
    return s or "deck"


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide.shapes.title.text = title
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: List[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    if not bullets:
        p = tf.paragraphs[0]
        p.text = "（预留）"
        return
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.level = 0


def build_lesson_deck(lesson: Dict[str, Any], chapter_title: str) -> Presentation:
    title = f"第{lesson['number']}讲 | {lesson['title']}"
    practical = lesson.get('practical', False)
    topics: List[str] = list(lesson.get('topics', []))
    code_paths: List[str] = list(lesson.get('code', []))

    prs = Presentation()

    # Title
    subtitle_bits = [chapter_title]
    if practical:
        subtitle_bits.append("实操")
    add_title_slide(prs, title, " · ".join(subtitle_bits))

    # Objectives (first 3 topics)
    add_bullet_slide(prs, "学习目标", topics[:3])

    # Key concepts (remaining topics)
    add_bullet_slide(prs, "核心概念", topics[3:6] if len(topics) > 3 else [])

    # Labs / Notebooks
    add_bullet_slide(prs, "实操与 Notebook", [str(p) for p in code_paths])

    # Lab guidance (common template)
    lab_tips = [
        "阅读 Notebook 注释，先跑通最小路径",
        "修改 1-2 处参数，观察行为变化",
        "在 Studio 中查看节点/边与状态流",
        "记录耗时/成本数据，准备讨论",
    ]
    add_bullet_slide(prs, "实验建议", lab_tips)

    # Summary
    add_bullet_slide(prs, "小结与下一讲", [
        "回顾本讲关键概念与实践",
        "梳理你的问题与改进点",
        "预习下一讲相关 Notebook",
    ])

    return prs


def build_chapter_deck(chapter: Dict[str, Any]) -> Presentation:
    prs = Presentation()
    chapter_title = chapter.get('chapter_title', '未命名章节')
    add_title_slide(prs, chapter_title, "教学大纲与要点")

    # Overview slide
    overview = [f"第{lesson['number']}讲：{lesson['title']}" for lesson in chapter.get('lessons', [])]
    add_bullet_slide(prs, "课程安排", overview)

    # Add one summary slide per lesson
    for lesson in chapter.get('lessons', []):
        title = f"第{lesson['number']}讲：{lesson['title']}"
        bullets = ["实操" if lesson.get('practical') else "理论", *(lesson.get('topics', [])[:4])]
        add_bullet_slide(prs, title, bullets)

    return prs


def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description='Build PPTX from syllabus.yaml')
    parser.add_argument('--syllabus', type=Path, default=Path('docs/syllabus.yaml'))
    parser.add_argument('--out', '--outdir', dest='outdir', type=Path, default=Path('slides'))
    parser.add_argument('--chapter', type=str, default='all', help='chapter id or "all"')
    parser.add_argument('--grouping', choices=['lesson', 'chapter'], default='lesson')
    parser.add_argument('--filter', nargs='*', default=[], help='filter lesson numbers (e.g., 1 2 3 15.1)')
    args = parser.parse_args(argv)

    data = load_syllabus(args.syllabus)
    chapters: List[Dict[str, Any]] = data.get('syllabus', [])

    args.outdir.mkdir(parents=True, exist_ok=True)

    selected_chapters = []
    if args.chapter == 'all':
        selected_chapters = chapters
    else:
        try:
            cid = int(float(args.chapter))
        except ValueError:
            print(f"Invalid --chapter: {args.chapter}", file=sys.stderr)
            return 2
        selected_chapters = [c for c in chapters if c.get('chapter_id') == cid]
        if not selected_chapters:
            print(f"No chapter matched id {cid}", file=sys.stderr)
            return 3

    lesson_filters = set(args.filter)

    built = 0
    if args.grouping == 'lesson':
        for chapter in selected_chapters:
            ctitle = chapter.get('chapter_title', '')
            for lesson in chapter.get('lessons', []):
                num_str = str(lesson.get('number'))
                if lesson_filters and num_str not in lesson_filters:
                    continue
                prs = build_lesson_deck(lesson, ctitle)
                fname = safe_filename(f"第{lesson['number']}讲_{lesson['title']}.pptx")
                outpath = args.outdir / fname
                prs.save(str(outpath))
                built += 1
    else:  # chapter grouping
        for chapter in selected_chapters:
            prs = build_chapter_deck(chapter)
            fname = safe_filename(f"{chapter.get('chapter_title', '未命名章节')}.pptx")
            outpath = args.outdir / fname
            prs.save(str(outpath))
            built += 1

    print(f"Built {built} deck(s) into {args.outdir}")
    return 0


if __name__ == '__main__':
    raise SystemExit(main())

